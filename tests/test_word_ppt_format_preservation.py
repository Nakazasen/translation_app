import concurrent.futures
from pathlib import Path

import pytest
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches as PptInches, Pt

from translation_app.core.file_handlers.powerpoint_handler import PowerPointHandler
from translation_app.core.file_handlers.word_handler import WordHandler


class FakeTranslationService:
    def __init__(self, translations=None):
        self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
        self.timeout = 5
        self.strategy = "ai"
        self.observer = None
        self.calls = []
        self.translations = translations or {}

    def set_runtime_observer(self, observer):
        self.observer = observer

    def clear_runtime_observer(self):
        self.observer = None

    def translate_long_text(self, text, src_lang, dest_lang):
        self.calls.append(text)
        if self.observer:
            self.observer("provider_call", {"provider": "fake"})
        return self.translations.get(text, f"{text}-{dest_lang}")

    def translate_text(self, text, src_lang, dest_lang):
        return self.translate_long_text(text, src_lang, dest_lang)


def _make_image(image_path: Path):
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow not installed")

    Image.new("RGB", (16, 16), color="blue").save(image_path)


def _run_word_translation(tmp_path, build_document, translations=None):
    input_path = tmp_path / "input.docx"
    output_path = tmp_path / "output.docx"
    service = FakeTranslationService(translations=translations)

    doc = Document()
    build_document(doc)
    doc.save(input_path)

    handler = WordHandler(service)
    try:
        stats = handler.translate(str(input_path), str(output_path), "en", "vi")
    finally:
        service.executor.shutdown(wait=True)

    return output_path, service, stats


def _run_ppt_translation(tmp_path, build_presentation, translations=None):
    input_path = tmp_path / "input.pptx"
    output_path = tmp_path / "output.pptx"
    service = FakeTranslationService(translations=translations)

    prs = Presentation()
    build_presentation(prs)
    prs.save(input_path)

    handler = PowerPointHandler(service)
    try:
        stats = handler.translate(str(input_path), str(output_path), "en", "vi")
    finally:
        service.executor.shutdown(wait=True)

    return output_path, service, stats


def test_word_preserves_run_formatting(tmp_path):
    def build_document(doc):
        p = doc.add_paragraph()
        p.style = doc.styles["Normal"]
        run1 = p.add_run("Hello ")
        run1.bold = True
        run2 = p.add_run("world")
        run2.italic = True
        run2.underline = True

    output_path, service, _ = _run_word_translation(tmp_path, build_document)

    result = Document(output_path)
    paragraph = result.paragraphs[0]
    assert len(paragraph.runs) == 2
    assert paragraph.runs[0].text == "Hello -vi"
    assert paragraph.runs[0].bold is True
    assert paragraph.runs[1].text == "world-vi"
    assert paragraph.runs[1].italic is True
    assert paragraph.runs[1].underline is True
    assert service.calls == ["Hello ", "world"]


def test_word_preserves_tables(tmp_path):
    def build_document(doc):
        table = doc.add_table(rows=2, cols=2)
        table.style = "Table Grid"
        table.cell(0, 0).text = "Alpha"
        table.cell(0, 1).text = "Beta"
        table.cell(1, 0).text = "Gamma"
        table.cell(1, 1).text = "Delta"

    output_path, service, _ = _run_word_translation(tmp_path, build_document)

    result = Document(output_path)
    table = result.tables[0]
    assert len(result.tables) == 1
    assert len(table.rows) == 2
    assert len(table.columns) == 2
    assert table.style.name == "Table Grid"
    assert table.cell(0, 0).text == "Alpha-vi"
    assert table.cell(0, 1).text == "Beta-vi"
    assert table.cell(1, 0).text == "Gamma-vi"
    assert table.cell(1, 1).text == "Delta-vi"
    assert service.calls == ["Alpha", "Beta", "Gamma", "Delta"]


def test_word_preserves_heading_and_bullets(tmp_path):
    def build_document(doc):
        heading = doc.add_paragraph(style="Heading 1")
        heading.add_run("Overview")
        bullet = doc.add_paragraph(style="List Bullet")
        bullet.add_run("Item one")

    output_path, service, _ = _run_word_translation(tmp_path, build_document)

    result = Document(output_path)
    assert result.paragraphs[0].style.name == "Heading 1"
    assert result.paragraphs[0].text == "Overview-vi"
    assert result.paragraphs[1].style.name == "List Bullet"
    assert result.paragraphs[1].text == "Item one-vi"
    assert service.calls == ["Overview", "Item one"]


def test_word_preserves_images(tmp_path):
    image_path = tmp_path / "tiny.png"
    _make_image(image_path)

    def build_document(doc):
        doc.add_paragraph("Caption")
        doc.add_picture(str(image_path), width=Inches(1))

    output_path, service, stats = _run_word_translation(tmp_path, build_document)

    result = Document(output_path)
    assert result.paragraphs[0].text == "Caption-vi"
    assert len(result.inline_shapes) == 1
    assert stats["images_processed"] == 0
    assert stats["images_skipped"] == 0
    assert service.calls == ["Caption"]


def test_word_skips_urls_emails_or_fields_if_supported(tmp_path):
    def build_document(doc):
        p = doc.add_paragraph()
        p.add_run("https://example.com")
        p.add_run(" ")
        p.add_run("contact@example.com")
        p.add_run(" ")
        p.add_run("C:\\temp\\file.txt")
        p.add_run(" ")
        p.add_run("Translate me")

    output_path, service, _ = _run_word_translation(tmp_path, build_document)

    result = Document(output_path)
    assert result.paragraphs[0].runs[0].text == "https://example.com"
    assert result.paragraphs[0].runs[2].text == "contact@example.com"
    assert result.paragraphs[0].runs[4].text == "C:\\temp\\file.txt"
    assert result.paragraphs[0].runs[6].text == "Translate me-vi"
    assert service.calls == ["Translate me"]


def test_ppt_preserves_shape_position_size(tmp_path):
    def build_presentation(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(3), PptInches(1))
        shape.text_frame.text = "Hello slide"

    output_path, service, _ = _run_ppt_translation(tmp_path, build_presentation)

    result = Presentation(output_path)
    shape = result.slides[0].shapes[0]
    assert shape.left == PptInches(1)
    assert shape.top == PptInches(1.5)
    assert shape.width == PptInches(3)
    assert shape.height == PptInches(1)
    assert shape.text == "Hello slide-vi"
    assert service.calls == ["Hello slide"]


def test_ppt_preserves_run_formatting(tmp_path):
    def build_presentation(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(PptInches(1), PptInches(1), PptInches(4), PptInches(1.5))
        paragraph = shape.text_frame.paragraphs[0]
        run1 = paragraph.add_run()
        run1.text = "Hello "
        run1.font.bold = True
        run1.font.size = Pt(20)
        run2 = paragraph.add_run()
        run2.text = "world"
        run2.font.italic = True
        run2.font.size = Pt(24)

    output_path, service, _ = _run_ppt_translation(tmp_path, build_presentation)

    result = Presentation(output_path)
    paragraph = result.slides[0].shapes[0].text_frame.paragraphs[0]
    assert len(paragraph.runs) >= 2
    assert paragraph.runs[0].text == "Hello -vi"
    assert paragraph.runs[0].font.bold is True
    assert paragraph.runs[0].font.size == Pt(20)
    assert paragraph.runs[1].text == "world-vi"
    assert paragraph.runs[1].font.italic is True
    assert paragraph.runs[1].font.size == Pt(24)
    assert service.calls == ["Hello ", "world"]


def test_ppt_preserves_tables(tmp_path):
    def build_presentation(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 2, PptInches(1), PptInches(1), PptInches(4), PptInches(1.5))
        table = table_shape.table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(1, 0).text = "C"
        table.cell(1, 1).text = "D"

    output_path, service, _ = _run_ppt_translation(tmp_path, build_presentation)

    result = Presentation(output_path)
    table = result.slides[0].shapes[0].table
    assert len(table.rows) == 2
    assert len(table.columns) == 2
    assert table.cell(0, 0).text == "A-vi"
    assert table.cell(0, 1).text == "B-vi"
    assert table.cell(1, 0).text == "C-vi"
    assert table.cell(1, 1).text == "D-vi"
    assert service.calls == ["A", "B", "C", "D"]


def test_ppt_preserves_images(tmp_path):
    image_path = tmp_path / "tiny.png"
    _make_image(image_path)

    def build_presentation(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), PptInches(1), PptInches(1), width=PptInches(1), height=PptInches(1))
        textbox = slide.shapes.add_textbox(PptInches(2.5), PptInches(1), PptInches(2), PptInches(1))
        textbox.text_frame.text = "Caption"

    output_path, service, stats = _run_ppt_translation(tmp_path, build_presentation)

    result = Presentation(output_path)
    shapes = result.slides[0].shapes
    picture_count = sum(1 for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert picture_count == 1
    assert shapes[1].text == "Caption-vi"
    assert stats["images_processed"] == 0
    assert stats["images_skipped"] == 0
    assert service.calls == ["Caption"]


def test_ppt_preserves_notes_if_supported(tmp_path):
    def build_presentation(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(PptInches(1), PptInches(1), PptInches(2), PptInches(1))
        textbox.text_frame.text = "Main"
        notes = slide.notes_slide
        for shape in notes.shapes:
            if getattr(shape, "has_text_frame", False) and "Notes" in shape.name:
                shape.text_frame.text = "Speaker notes"

    output_path, service, _ = _run_ppt_translation(tmp_path, build_presentation)

    result = Presentation(output_path)
    notes = result.slides[0].notes_slide
    notes_text = ""
    for shape in notes.shapes:
        if getattr(shape, "has_text_frame", False) and "Notes" in shape.name:
            notes_text = shape.text_frame.text
            break
    assert result.slides[0].shapes[0].text == "Main-vi"
    assert notes_text == "Speaker notes-vi"
    assert service.calls == ["Main", "Speaker notes"]
