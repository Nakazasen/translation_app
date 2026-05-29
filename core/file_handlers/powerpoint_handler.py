"""
PowerPoint file handler for translation.
"""
import re
from typing import Any, Dict

from pptx import Presentation

from translation_app.core.file_translation_control import FileTranslationInterrupted, FileTranslationStopRequested
from translation_app.core.ocr_handler import get_ocr_handler
from translation_app.core.translator import TranslationService
from translation_app.utils.error_handler import FileProcessingError
from translation_app.utils.logger import logger


class PowerPointHandler:
    """Handler for PowerPoint file translation."""

    _URL_RE = re.compile(r"^(?:https?://|www\.)\S+$", re.IGNORECASE)
    _EMAIL_RE = re.compile(r"^[^@\s]+@[^@\s]+\.[^@\s]+$")
    _FILE_PATH_RE = re.compile(r"^(?:[A-Za-z]:\\|\\\\|/)[^\r\n]+$")
    _FIELD_HINT_RE = re.compile(r"\b(?:HYPERLINK|MERGEFIELD|PAGE|NUMPAGES|TOC)\b", re.IGNORECASE)
    _FLIGHT_CODE_RE = re.compile(r"^[A-Z]{1,4}\d{1,4}[A-Z]?$")
    _NUMERICISH_TOKEN_RE = re.compile(r"^(?=.*\d)[\d\s.,:/()%+\-]+$")
    _FRAGMENTED_RUN_MIN_COUNT = 2
    _FRAGMENTED_RUN_MIN_TOTAL_CHARS = 20

    def __init__(self, translation_service: TranslationService):
        self.translation_service = translation_service
        self.ocr_handler = get_ocr_handler()

    def translate(self, input_file: str, output_file: str, src_lang: str, dest_lang: str) -> Dict[str, Any]:
        """Translate a PPTX file while preserving core formatting."""
        prs = None
        try:
            logger.info(f"Starting PowerPoint translation: {input_file}")
            prs = Presentation(input_file)

            for slide in prs.slides:
                getattr(self.translation_service, "raise_if_file_translation_stopped", lambda: None)()
                self._translate_shapes_in_slide(slide, src_lang, dest_lang)
                self._translate_notes(slide, src_lang, dest_lang)

            logger.info("Skipping PowerPoint image OCR/write-back to preserve layout during hardening phase")
            logger.info("Skipping PowerPoint diagram/SmartArt deep translation in hardening phase")

            prs.save(output_file)
            logger.info(f"PowerPoint translation completed: {output_file}")

            return {
                "images_processed": 0,
                "images_skipped": 0,
                "diagrams_translated": 0,
            }
        except FileTranslationStopRequested as exc:
            partial_saved, save_error = self._save_partial_presentation(prs, output_file, exc.status)
            raise FileTranslationInterrupted(
                exc.status,
                output_file=output_file,
                partial_saved=partial_saved,
                save_error=save_error,
            ) from exc
        except Exception as exc:
            error_msg = f"Error translating PowerPoint file: {exc}"
            logger.error(error_msg)
            raise FileProcessingError(error_msg, original_error=exc) from exc

    def _translate_shapes_in_slide(self, slide, src_lang: str, dest_lang: str) -> None:
        for shape in slide.shapes:
            getattr(self.translation_service, "raise_if_file_translation_stopped", lambda: None)()
            self._translate_single_shape(shape, src_lang, dest_lang)

    def _translate_single_shape(self, shape, src_lang: str, dest_lang: str) -> None:
        try:
            shape_type = shape.shape_type

            if shape_type == 6 and hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    self._translate_single_shape(sub_shape, src_lang, dest_lang)
                return

            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame is not None:
                            self._translate_text_frame(cell.text_frame, src_lang, dest_lang)
                return

            if hasattr(shape, "has_chart") and shape.has_chart:
                logger.info("Skipping chart translation to preserve chart structure during hardening phase")
                return

            if shape_type == 14 and not getattr(shape, "has_table", False):
                logger.info("Skipping SmartArt/diagram translation to preserve layout during hardening phase")
                return

            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                self._translate_text_frame(shape.text_frame, src_lang, dest_lang)
        except Exception as exc:
            logger.warning(f"Error processing PowerPoint shape: {exc}")

    def _translate_text_frame(self, text_frame, src_lang: str, dest_lang: str) -> None:
        for paragraph in text_frame.paragraphs:
            getattr(self.translation_service, "raise_if_file_translation_stopped", lambda: None)()
            eligible_runs = self._collect_translatable_runs(paragraph.runs)
            if self._should_translate_fragmented_runs_as_block(paragraph.runs, eligible_runs):
                self._translate_fragmented_paragraph(paragraph, src_lang, dest_lang)
                continue

            for run in eligible_runs:
                getattr(self.translation_service, "raise_if_file_translation_stopped", lambda: None)()
                original_text = run.text
                try:
                    translated_text = self.translation_service.translate_long_text(
                        original_text,
                        src_lang,
                        dest_lang,
                    )
                except Exception as exc:
                    logger.warning(f"Error translating PowerPoint run: {exc}")
                    continue
                if translated_text != original_text:
                    run.text = translated_text

    def _translate_notes(self, slide, src_lang: str, dest_lang: str) -> None:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            return

        if notes_slide is None:
            return

        notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
        if notes_text_frame is None:
            return

        self._translate_text_frame(notes_text_frame, src_lang, dest_lang)

    def _save_partial_presentation(self, prs, output_file: str, status: str) -> tuple[bool, Exception | None]:
        if prs is None:
            return False, None
        try:
            prs.save(output_file)
            logger.info(f"Saved partial PowerPoint output after file translation was {status}: {output_file}")
            return True, None
        except Exception as exc:
            logger.error(f"Failed to save partial PowerPoint output after {status}: {exc}")
            return False, exc

    def _should_skip_text(self, text: str) -> bool:
        stripped = (text or "").strip()
        if not stripped:
            return True
        if self._URL_RE.match(stripped):
            return True
        if self._EMAIL_RE.match(stripped):
            return True
        if self._FILE_PATH_RE.match(stripped):
            return True
        if self._FIELD_HINT_RE.search(stripped):
            return True
        if self._is_nonlinguistic_ascii_token(stripped):
            return True
        return False

    def _is_nonlinguistic_ascii_token(self, stripped: str) -> bool:
        if self._FLIGHT_CODE_RE.match(stripped):
            return True
        if self._NUMERICISH_TOKEN_RE.match(stripped):
            return True
        if not stripped.isascii() or not any(ch.isdigit() for ch in stripped):
            return False
        words = re.findall(r"[A-Za-z]+", stripped)
        return bool(words) and all(len(word) <= 2 for word in words)

    def _collect_translatable_runs(self, runs) -> list:
        return [run for run in runs if not self._should_skip_text(run.text)]

    def _should_translate_fragmented_runs_as_block(self, all_runs, eligible_runs) -> bool:
        if len(eligible_runs) < self._FRAGMENTED_RUN_MIN_COUNT:
            return False

        for run in all_runs:
            run_text = run.text or ""
            if not run_text.strip():
                continue
            if self._should_skip_text(run_text):
                return False

        total_chars = sum(len((run.text or "").strip()) for run in eligible_runs)
        return total_chars >= self._FRAGMENTED_RUN_MIN_TOTAL_CHARS

    def _translate_fragmented_paragraph(self, paragraph, src_lang: str, dest_lang: str) -> None:
        original_text = paragraph.text
        try:
            translated_text = self.translation_service.translate_long_text(
                original_text,
                src_lang,
                dest_lang,
            )
        except Exception as exc:
            logger.warning(f"Error translating fragmented PowerPoint paragraph: {exc}")
            return

        if translated_text == original_text:
            return

        first_run = None
        for run in paragraph.runs:
            if first_run is None:
                first_run = run
                run.text = translated_text
                continue
            run.text = ""
